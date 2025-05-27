import os
import logging
import re
import time
import io
import writer as wf
import writer.ai
import pandas as pd
from typing import Literal, Any
from prompts import (
    get_outline_prompt,
    get_visual_outline_prompt,
    get_briefing_info_prompt,
    get_map_messages_prompt,
    get_speaker_notes_prompt,
)

# Try to import slide intelligence, but make it optional
try:
    from slide_intelligence import (
        analyze_presentation_intelligence,
        format_intelligent_notes,
        SlideType,
        SlideAnalyzer
    )
    INTELLIGENCE_SUPPORT = True
except ImportError:
    INTELLIGENCE_SUPPORT = False
    logger.warning("slide_intelligence module not found. Advanced intelligence features disabled.")

# Try to import language utilities
try:
    from language_utils import detect_and_adapt_prompts, enhance_prompt_with_language
    LANGUAGE_SUPPORT = True
except ImportError:
    LANGUAGE_SUPPORT = False
    logger.warning("language_utils module not found. Multi-language support disabled.")

try:
    from pptx import Presentation
    PPTX_SUPPORT = True
except ImportError:
    PPTX_SUPPORT = False

try:
    import docx
    from docx.shared import Inches, Pt
    DOCX_SUPPORT = True
except ImportError:
    DOCX_SUPPORT = False


logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# Enable debug logging for troubleshooting
# Uncomment the following line to see debug messages
logger.setLevel(logging.DEBUG)

if "WRITER_API_KEY" in os.environ:
    wf.api_key = os.getenv("WRITER_API_KEY")
else:
    logger.warning("WRITER_API_KEY environment variable not set. AI features will not work.")

# --- Custom Exception for Stopping ---
class GenerationStoppedError(Exception):
    """Custom exception for stopping the generation process."""
    pass

# --- Helper Functions ---
def _extract_text_from_file(file_path: str) -> str:
    """Extracts text from PDF or PPTX files."""
    logger.info(f"Attempting to extract text from: {file_path}")
    file_extension = os.path.splitext(file_path)[1].lower()
    base_filename = os.path.basename(file_path)
    sanitized_filename = base_filename.encode('ascii', 'ignore').decode('ascii')

    try:
        if file_extension == ".pdf":
            logger.info("Processing as PDF using writer.ai.tools.parse_pdf")
            with open(file_path, "rb") as f:
                file_data = f.read()

            if not wf.api_key:
                raise ConnectionError("WRITER_API_KEY is not set, cannot upload file.")

            uploaded_file = writer.ai.upload_file(
                data=file_data,
                type="application/pdf",
                name=sanitized_filename
            )
            file_id = uploaded_file.id
            logger.info(f"Uploaded PDF {sanitized_filename} with ID: {file_id}")

            extracted_text = writer.ai.tools.parse_pdf(file_id_or_file=file_id, format="text")
            logger.info(f"Extracted text from PDF {sanitized_filename}")
            try:
                writer.ai.delete_file(file_id)
                logger.info(f"Deleted temporary file {file_id}")
            except Exception as del_e:
                logger.warning(f"Could not delete temporary file {file_id}: {del_e}")
            return extracted_text

        elif file_extension == ".pptx":
            if not PPTX_SUPPORT:
                 logger.warning(f"Skipping PPTX text extraction for {base_filename} as python-pptx is not installed.")
                 return f"[Text extraction skipped for PPTX: {base_filename} - python-pptx not installed]"

            logger.info(f"Processing as PPTX locally using python-pptx")
            presentation = Presentation(file_path)
            full_text = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            full_text.append(text)
            extracted_text = "\n".join(full_text)
            logger.info(f"Extracted text from PPTX {base_filename}")
            return extracted_text

        else:
            logger.warning(f"Unsupported file type for text extraction: {base_filename}. Returning placeholder.")
            return f"[Unsupported file type: {base_filename}]"

    except Exception as e:
        logger.error(f"Error extracting text from {base_filename}: {e}", exc_info=True)
        raise RuntimeError(f"Failed to extract text from {base_filename}") from e

def _call_ai_model(prompt: str, model: str = "palmyra-x-004", temperature: float = 0.0) -> str:
    """Calls the Writer AI completion endpoint with a specific temperature."""
    if not wf.api_key:
        raise ConnectionError("WRITER_API_KEY is not set, cannot call AI model.")
    try:
        config = {"model": model, "temperature": temperature}
        response = writer.ai.complete(prompt, config=config)
        return response.strip()
    except Exception as e:
        logger.error(f"Error calling AI model {model} with temperature {temperature}: {e}")
        raise

def _analyze_visuals(file_path: str, prompt: str) -> str:
    """Analyzes visuals (or text file content) in a presentation using Palmyra Vision."""
    if not wf.api_key:
        raise ConnectionError("WRITER_API_KEY is not set, cannot analyze visuals.")

    file_extension = os.path.splitext(file_path)[1].lower()
    base_filename = os.path.basename(file_path)
    sanitized_filename = base_filename.encode('ascii', 'ignore').decode('ascii')
    
    # Extended list of supported formats for better visual analysis
    supported_vision_types = [
        ".pdf", ".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff", ".svg", ".txt"
    ]
    
    if file_extension not in supported_vision_types:
        logger.warning(f"File type {file_extension} not directly supported by Palmyra Vision. Attempting text extraction fallback.")
        # For unsupported types, try text extraction if it's a presentation file
        if file_extension in [".pptx", ".ppt"]:
            try:
                extracted_text = _extract_text_from_file(file_path)
                if extracted_text:
                    # Save as temporary text file for vision analysis
                    temp_text_path = os.path.join("data", f"temp_visual_{int(time.time())}.txt")
                    with open(temp_text_path, "w", encoding="utf-8") as f:
                        f.write(extracted_text)
                    
                    # Recursively call with text file
                    result = _analyze_visuals(temp_text_path, prompt)
                    
                    # Clean up temp file
                    if os.path.exists(temp_text_path):
                        os.remove(temp_text_path)
                    
                    return result
            except Exception as e:
                logger.error(f"Failed to extract text for visual analysis: {e}")
        
        return "[Visual analysis skipped: Unsupported file type]"

    file_id = None
    try:
        with open(file_path, "rb") as f:
            file_data = f.read()

        if file_extension == ".pdf":
            content_type = "application/pdf"
        elif file_extension == ".txt":
            content_type = "text/plain"
        elif file_extension in [".svg"]:
            content_type = "image/svg+xml"
        else:
            import mimetypes
            content_type, _ = mimetypes.guess_type(file_path)
            if not content_type:
                # Default content types for common image formats
                content_type_map = {
                    ".png": "image/png",
                    ".jpg": "image/jpeg",
                    ".jpeg": "image/jpeg",
                    ".gif": "image/gif",
                    ".webp": "image/webp",
                    ".bmp": "image/bmp",
                    ".tiff": "image/tiff"
                }
                content_type = content_type_map.get(file_extension, "application/octet-stream")

        uploaded_file = writer.ai.upload_file(
            data=file_data,
            type=content_type,
            name=sanitized_filename
        )
        file_id = uploaded_file.id
        logger.info(f"Uploaded file for visual analysis {sanitized_filename} with ID: {file_id}")

        # All prompts should now use {{InputDocument}}
        placeholder_found = "{{InputDocument}}" in prompt
        logger.info(f"Vision API check - Placeholder found: {placeholder_found}, Prompt length: {len(prompt)}")
        
        if placeholder_found:
            variables = [{"name": "InputDocument", "file_id": file_id}]
        else:
            logger.warning("Prompt does not contain {{InputDocument}} placeholder. Using empty variables array.")
            logger.info(f"First 200 chars of prompt: {prompt[:200]}")
            variables = []  # Empty array, not None

        client = writer.ai.WriterAIManager.acquire_client()
        
        # Log the prompt for debugging (first 200 chars)
        logger.debug(f"Vision API prompt preview: {prompt[:200]}...")
        logger.info(f"Vision API variables count: {len(variables)}")
        
        try:
            # Only pass variables parameter if we have actual variables
            if variables:
                response = client.vision.analyze(
                    model="palmyra-vision",
                    prompt=prompt,
                    variables=variables
                )
            else:
                # Omit variables parameter entirely if empty
                response = client.vision.analyze(
                    model="palmyra-vision",
                    prompt=prompt
                )
            logger.info(f"Visual analysis completed for {sanitized_filename}")
            
            # Check if we got a valid response
            if hasattr(response, 'data') and response.data:
                result = response.data
            else:
                logger.warning(f"Vision API returned empty or invalid response for {sanitized_filename}")
                result = "[Visual analysis returned no content]"
        except Exception as api_error:
            logger.error(f"Vision API error: {api_error}")
            # Try to provide more helpful error message
            if "variables" in str(api_error):
                logger.error("Vision API variables issue. Current variables: " + str(variables))
            raise
        logger.info(f"Visual analysis completed for {sanitized_filename}")

        try:
            writer.ai.delete_file(file_id)
            logger.info(f"Deleted temporary file {file_id} after visual analysis")
        except Exception as del_e:
            logger.warning(f"Could not delete temporary file {file_id} after visual analysis: {del_e}")

        return result
    except Exception as e:
        logger.error(f"Error analyzing visuals in {base_filename}: {e}", exc_info=True)
        if file_id:
             try:
                 writer.ai.delete_file(file_id)
                 logger.info(f"Deleted temporary file {file_id} after analysis error")
             except Exception as del_e:
                 logger.warning(f"Could not delete temporary file {file_id} after analysis error: {del_e}")
        return f"[Error during visual analysis: {type(e).__name__}]"

def _add_formatted_text_to_paragraph(paragraph: Any, text: str):
    """Adds text to a paragraph, applying bold and italic formatting based on Markdown."""
    parts = re.split(r'(\*\*.*?\*\*|[*_].*?[*_])', text)

    for part in parts:
        if not part:
            continue

        run = paragraph.add_run()
        if part.startswith('**') and part.endswith('**'):
            run.text = part[2:-2]
            run.bold = True
        elif (part.startswith('*') and part.endswith('*')) or \
             (part.startswith('_') and part.endswith('_')):
            run.text = part[1:-1]
            run.italic = True
        else:
            run.text = part

def _generate_docx_bytes(content: str) -> bytes:
    """Creates a DOCX file in memory from markdown-like text content and returns its bytes."""
    if not DOCX_SUPPORT:
        logger.error("python-docx library not found. Cannot create DOCX file.")
        raise ImportError("python-docx is required for DOCX generation. Please install it.")

    document = docx.Document()

    for line in content.split('\n'):
        stripped_line = line.strip()

        heading_match = re.match(r"^(#+)\s+(.*)", stripped_line)
        if heading_match:
            level = len(heading_match.group(1))
            text_content = heading_match.group(2)
            level = max(1, min(level, 9))
            heading = document.add_heading(level=level)
            _add_formatted_text_to_paragraph(heading, text_content)
        elif stripped_line.startswith("Slide ") and ":" in stripped_line:
             try:
                 heading_text = stripped_line.split(":", 1)[1].strip()
                 heading = document.add_heading(level=1)
                 _add_formatted_text_to_paragraph(heading, heading_text)
             except IndexError:
                 heading = document.add_heading(level=1)
                 _add_formatted_text_to_paragraph(heading, stripped_line)
        elif stripped_line.startswith("* ") or stripped_line.startswith("- "):
            text_content = stripped_line[2:]
            if text_content:
                 paragraph = document.add_paragraph(style='List Bullet')
                 _add_formatted_text_to_paragraph(paragraph, text_content)
            else:
                 document.add_paragraph(style='List Bullet')
        elif not stripped_line:
             document.add_paragraph()
        else:
            paragraph = document.add_paragraph()
            _add_formatted_text_to_paragraph(paragraph, stripped_line)

    buffer = io.BytesIO()
    try:
        document.save(buffer)
        logger.info("Generated DOCX content in memory.")
        buffer.seek(0)
        return buffer.getvalue()
    except Exception as e:
        logger.error(f"Failed to save DOCX to memory buffer: {e}")
        raise

def _download_content_as_docx(state, content_key: str):
    """Helper to generate and download content as DOCX."""
    logger.info(f"Attempting to download content for key: {content_key}")
    
    # Check if results exist and has the content key
    if "results" not in state or content_key not in state["results"]:
        message = f"!Results not available for download."
        state["processing_message"] = message
        logger.warning(f"Results key missing in state: {list(state.keys())}")
        return
        
    content = state["results"].get(content_key, "")
    
    if not content:
        message = f"!Content for '{content_key}' not available for download."
        state["processing_message"] = message
        logger.warning(f"Content empty for key: {content_key}")
        return

    # Ensure content is string
    if not isinstance(content, str):
        content = str(content)

    try:
        filename_docx = f"{content_key}_notes.docx"
        docx_bytes = _generate_docx_bytes(content)
        state.file_download(
            wf.pack_bytes(docx_bytes, "application/vnd.openxmlformats-officedocument.wordprocessingml.document"),
            filename_docx
        )
        state["processing_message"] = f"+Download initiated for {filename_docx}."
        logger.info(f"Download successful for {filename_docx}")
    except Exception as e:
        message = f"-Error preparing DOCX download for '{content_key}': {e}"
        state["processing_message"] = message
        logger.error(message, exc_info=True)

def _get_input_text(state, input_type: Literal["deck", "briefing"]) -> str:
    """Gets the text content based on the selected input method."""
    if input_type == "deck":
        method = state["deck_input_method"]
        if method == "Upload File":
            if state["deck_file"]["path"]:
                if state["deck_file"]["text_content"] is None:
                     state["processing_message"] = "%Extracting text from presentation deck..."
                     state["deck_file"]["text_content"] = _extract_text_from_file(state["deck_file"]["path"])
                return state["deck_file"]["text_content"] or ""
            else:
                return ""
        elif method == "Enter Text":
            return state["deck_file"]["text_content"] or ""
        else:
            return ""
    elif input_type == "briefing":
        method = state["briefing_input_method"]
        if method == "Upload File":
            briefing_texts = []
            for i, briefing_file in enumerate(state["briefing_files"]):
                if briefing_file["path"]:
                    if briefing_file["text_content"] is None:
                        state["processing_message"] = f"%Extracting text from briefing document {i+1}/{len(state['briefing_files'])}..."
                        briefing_file["text_content"] = _extract_text_from_file(briefing_file["path"])
                    briefing_texts.append(briefing_file["text_content"] or "")
            return "\n\n---\n\n".join(briefing_texts)
        elif method == "Enter Text":
            return state["briefing_manual_text"] or ""
        else:
            return ""
    return ""

def _update_generate_button_state(state):
    """Updates the disabled state of the generate button based on inputs."""
    deck_ready = False
    if state["deck_input_method"] == "Upload File":
        deck_ready = bool(state["deck_file"]["path"])
    elif state["deck_input_method"] == "Enter Text":
        deck_ready = bool(state["deck_file"]["text_content"])

    # Briefing is now optional - button is enabled if deck is ready
    state["ui_controls"]["generate_disabled"] = "yes" if (state["is_generating"] or not deck_ready) else "no"


# --- State Initialization ---
os.makedirs("data", exist_ok=True)
placeholder_data = {'Description': [], 'Label': []}
initial_df = pd.DataFrame(placeholder_data)

initial_state = wf.init_state({
    "app_title": "Speaker Notes Generator",
    "logo_path": "static/writer_logo.svg",
    "deck_file": {"name": None, "path": None, "id": None, "text_content": None},
    "briefing_files": [],
    "settings": {
        "timing": "30 Minutes",
        "style": "Informative",
        "verbosity": "Standard"  # New setting for verbosity control
    },
    "metrics": {"new_features": 0, "caveats": 0, "fixed_issues": 0, "total": 0},
    "processing_step": "idle",
    "processing_message": "Upload your presentation deck and briefing document(s).",
    "results": {
        "outline": "",
        "visuals": "",
        "briefing_info": "",
        "mapping": "",
        "speaker_notes": "",
    },
    "ui_controls": {
        "generate_disabled": "yes",
        "download_disabled": "yes",
        "show_deck_upload": True,
        "show_deck_text_input": False,
        "show_briefing_upload": True,
        "show_briefing_text_input": False,
    },
    "deck_input_method": "Upload File",
    "briefing_input_method": "Upload File",
    "deck_text_input": "",
    "briefing_text_input": "",
    "briefing_manual_text": "",
    "is_generating": False,
    "stop_requested": False,
})
initial_state.import_stylesheet("custom_styles", "/static/custom.css")

# --- Event Handlers ---
def handle_deck_upload(state, payload):
    """Handles the presentation deck upload."""
    if not payload:
        state["processing_message"] = "-No deck file uploaded."
        return
    try:
        uploaded_file = payload[0]
        name = uploaded_file.get("name")
        file_data = uploaded_file.get("data")

        temp_path = os.path.join("data", f"deck_{name}")
        with open(temp_path, "wb") as file_handle:
            file_handle.write(file_data)

        if state["deck_file"] is None: state["deck_file"] = {}
        state["deck_file"]["name"] = name
        state["deck_file"]["path"] = temp_path
        state["deck_file"]["id"] = None
        state["deck_file"]["text_content"] = None
        state["processing_message"] = f"+Deck '{name}' uploaded."
        _update_generate_button_state(state)
        logger.info(f"Deck file uploaded: {name}")
    except Exception as e:
        state["processing_message"] = f"-Error uploading deck: {e}"
        logger.error(f"Error in handle_deck_upload: {e}", exc_info=True)

def handle_briefing_upload(state, payload):
    """Handles briefing document uploads."""
    if not payload:
        state["processing_message"] = "-No briefing file(s) uploaded."
        return
    try:
        new_files_info = []
        for uploaded_file in payload:
            name = uploaded_file.get("name")
            file_data = uploaded_file.get("data")
            temp_path = os.path.join("data", f"briefing_{name}")
            with open(temp_path, "wb") as file_handle:
                file_handle.write(file_data)
            new_files_info.append({"name": name, "path": temp_path, "id": None, "text_content": None})
            logger.info(f"Briefing file uploaded: {name}")

        if state["briefing_files"] is None: state["briefing_files"] = []
        state["briefing_files"] += new_files_info
        state["briefing_manual_text"] = ""
        state["processing_message"] = f"+{len(payload)} briefing file(s) added."
        _update_generate_button_state(state)

    except Exception as e:
        state["processing_message"] = f"-Error uploading briefing(s): {e}"
        logger.error(f"Error in handle_briefing_upload: {e}", exc_info=True)

def handle_deck_method_change(state, payload):
    state["deck_input_method"] = payload
    state["ui_controls"]["show_deck_upload"] = (payload == "Upload File")
    state["ui_controls"]["show_deck_text_input"] = (payload == "Enter Text")
    if payload == "Upload File":
        state["deck_file"]["text_content"] = None
    else:
        state["deck_file"]["path"] = None
        state["deck_file"]["name"] = None
        state["deck_file"]["id"] = None
    _update_generate_button_state(state)

def handle_briefing_method_change(state, payload):
    state["briefing_input_method"] = payload
    state["ui_controls"]["show_briefing_upload"] = (payload == "Upload File")
    state["ui_controls"]["show_briefing_text_input"] = (payload == "Enter Text")
    if payload == "Upload File":
        state["briefing_manual_text"] = ""
    else:
        state["briefing_files"] = []
    _update_generate_button_state(state)

def handle_save_deck_text(state):
    state["deck_file"]["text_content"] = state["deck_text_input"]
    state["deck_file"]["path"] = None
    state["deck_file"]["name"] = "Manual Text Input"
    state["deck_file"]["id"] = None
    state["processing_message"] = "+Deck text saved."
    _update_generate_button_state(state)

def handle_save_briefing_text(state):
    state["briefing_manual_text"] = state["briefing_text_input"]
    state["briefing_files"] = []
    state["processing_message"] = "+Briefing text saved."
    _update_generate_button_state(state)

def handle_stop_generate(state):
    """Sets the stop flag to interrupt the generation process."""
    if state["is_generating"]:
        state["stop_requested"] = True
        state["processing_message"] = "!Stopping generation..."
        logger.info("Stop generation requested by user.")
    else:
        state["processing_message"] = "!Generation is not currently running."

def handle_generate(state):
    """Orchestrates the AI processing pipeline with verbosity control."""
    deck_ready = False
    if state["deck_input_method"] == "Upload File":
        deck_ready = bool(state["deck_file"]["path"])
    elif state["deck_input_method"] == "Enter Text":
        deck_ready = bool(state["deck_file"]["text_content"])

    briefing_ready = False
    if state["briefing_input_method"] == "Upload File":
        briefing_ready = len(state["briefing_files"]) > 0
    elif state["briefing_input_method"] == "Enter Text":
        briefing_ready = bool(state["briefing_manual_text"])

    if not deck_ready or not briefing_ready:
        state["processing_message"] = "!Please provide input for both the deck and briefing documents."
        return

    if not wf.api_key:
         state["processing_message"] = "-WRITER_API_KEY is not set. Cannot proceed with AI generation."
         logger.error("WRITER_API_KEY not set.")
         return

    state["is_generating"] = True
    state["stop_requested"] = False
    state["ui_controls"]["generate_disabled"] = "yes"
    state["ui_controls"]["download_disabled"] = "yes"
    state["processing_step"] = "idle"
    state["results"] = {}
    
    # Get verbosity level
    verbosity = state["settings"]["verbosity"]
    timing = state["settings"]["timing"]
    style = state["settings"]["style"]

    try:
        # Step 0: Get Text Content
        state["processing_step"] = "text_extraction"
        state["processing_message"] = "%Getting text content..."
        deck_text = _get_input_text(state, "deck")
        combined_briefing_text = _get_input_text(state, "briefing")

        if not deck_text:
             raise ValueError("Presentation deck content is empty.")
        if not combined_briefing_text:
             raise ValueError("Briefing document content is empty.")

        # Detect language if supported
        if LANGUAGE_SUPPORT:
            try:
                language_info = detect_and_adapt_prompts(deck_text, combined_briefing_text)
                if language_info and language_info.get('language_code') != 'en':
                    state["processing_message"] = f"%Detected {language_info['language_name']} content. Generating notes in {language_info['language_name']}..."
                    logger.info(f"Language detected: {language_info['language_name']} ({language_info['language_code']})")
            except Exception as e:
                logger.warning(f"Language detection failed: {e}")
                language_info = None

        if state["stop_requested"]: raise GenerationStoppedError()

        # Step 1: Extract Outline with verbosity control and language support
        state["processing_step"] = "step1"
        state["processing_message"] = "%Generating presentation outline..."
        
        # Get the appropriate prompt based on verbosity
        outline_prompt = get_outline_prompt(verbosity)
        
        # Add language instructions if detected
        if language_info:
            try:
                outline_prompt = enhance_prompt_with_language(outline_prompt, language_info)
            except Exception as e:
                logger.warning(f"Failed to enhance prompt with language: {e}")
                # Continue with original prompt
        
        # Use regular AI model for outline generation
        logger.info("Generating outline using text analysis")
        full_prompt = outline_prompt + "\n" + deck_text
        outline = _call_ai_model(full_prompt, temperature=0.0)
        
        state["results"]["outline"] = outline
        logger.info("Step 1 (Outline) completed.")

        if state["stop_requested"]: raise GenerationStoppedError()

        # Step 2: Analyze Visuals with verbosity control and language support
        state["processing_step"] = "step2"
        if state["deck_input_method"] == "Upload File" and state["deck_file"]["path"]:
            state["processing_message"] = "%Analyzing presentation visuals..."
            visual_prompt = get_visual_outline_prompt(verbosity)
            
            # Add language instructions if detected
            if language_info:
                visual_prompt = enhance_prompt_with_language(visual_prompt, language_info)
                
            visual_analysis = _analyze_visuals(state["deck_file"]["path"], visual_prompt)
            state["results"]["visuals"] = visual_analysis
            logger.info("Step 2 (Visuals) completed.")
        else:
            state["results"]["visuals"] = "[Visual analysis skipped: Text input provided for deck]"
            logger.info("Step 2 (Visuals) skipped as text input was used for deck.")

        if state["stop_requested"]: raise GenerationStoppedError()

        # Step 3: Extract Briefing Info (if briefing provided) with language support
        state["processing_step"] = "step3"
        logger.info(f"Step 3 starting. has_briefing={has_briefing}, briefing text length={len(combined_briefing_text)}")
        
        if has_briefing:
            state["processing_message"] = "%Extracting key information from briefing..."
            briefing_prompt = get_briefing_info_prompt(verbosity)
            
            # Add language instructions if detected
            if language_info:
                briefing_prompt = enhance_prompt_with_language(briefing_prompt, language_info)
                
            prompt3 = briefing_prompt.replace("{{Briefing(s)}}", combined_briefing_text)
            briefing_info = _call_ai_model(prompt3, temperature=0.0)
            state["results"]["briefing_info"] = briefing_info
            logger.info("Step 3 (Briefing Info) completed.")
        else:
            # Generate context from deck content when no briefing is provided
            state["processing_message"] = "%Generating context from presentation content..."
            context_prompt = f"""
Analyze the presentation content and extract key context that would typically come from a briefing document.

Based on the presentation outline and visuals:
{outline}

{state["results"]["visuals"] or ""}

Extract and infer:
• Key Business Context: What situation or challenge is this presentation addressing?
• Target Audience: Who is this presentation likely for?
• Main Objectives: What are the apparent goals of this presentation?
• Key Themes: What are the recurring topics or messages?
• Implied Recommendations: What actions or decisions does the presentation seem to advocate?

Present findings in structured bullet points for easy reference.
"""
            # Add language instructions if detected
            if language_info:
                context_prompt = enhance_prompt_with_language(context_prompt, language_info)
                
            briefing_info = _call_ai_model(context_prompt, temperature=0.0)
            state["results"]["briefing_info"] = briefing_info
            logger.info("Step 3 (Context Generation from Deck) completed.")

        if state["stop_requested"]: raise GenerationStoppedError()

        # Step 4: Map Messages with verbosity control and language support
        state["processing_step"] = "step4"
        state["processing_message"] = "%Mapping messages to slides..."
        mapping_prompt = get_map_messages_prompt(verbosity)
        
        # Add language instructions if detected
        if language_info:
            mapping_prompt = enhance_prompt_with_language(mapping_prompt, language_info)
            
        prompt4 = mapping_prompt.replace("{{Presentation Outline}}", outline)
        prompt4 = prompt4.replace("{{Visual Presentation Outline}}", state["results"]["visuals"] or "")
        prompt4 = prompt4.replace("{{Key Information from the Briefing}}", briefing_info)
        mapping = _call_ai_model(prompt4, temperature=0.0)
        state["results"]["mapping"] = mapping
        logger.info("Step 4 (Mapping) completed.")

        if state["stop_requested"]: raise GenerationStoppedError()

        # Step 5: Generate Speaker Notes with optional intelligence enhancement and language support
        state["processing_step"] = "step5"
        state["processing_message"] = "%Generating speaker notes..."
        
        # Analyze presentation intelligence if module is available
        presentation_intel = {}
        if INTELLIGENCE_SUPPORT and verbosity != "Detailed":
            try:
                presentation_intel = analyze_presentation_intelligence(outline, state["results"]["visuals"] or "")
                logger.info(f"Analyzed {len(presentation_intel)} slides for intelligence")
            except Exception as e:
                logger.warning(f"Intelligence analysis failed: {e}")
        
        # Generate base speaker notes
        speaker_notes_prompt = get_speaker_notes_prompt(verbosity, timing, style)
        
        # Add language instructions if detected
        if language_info:
            speaker_notes_prompt = enhance_prompt_with_language(speaker_notes_prompt, language_info)
            
        prompt5 = speaker_notes_prompt.replace("{{Map Messages to Each Slide}}", mapping)
        prompt5 = prompt5.replace("{{Presentation Outline}}", outline)
        prompt5 = prompt5.replace("{{Visual Presentation Outline}}", state["results"]["visuals"] or "")
        speaker_notes_text = _call_ai_model(prompt5, temperature=0.0)
        
        # Enhance notes with slide intelligence if available and not in Detailed mode
        if INTELLIGENCE_SUPPORT and presentation_intel and verbosity != "Detailed":
            try:
                enhanced_notes = []
                notes_by_slide = re.split(r'(Slide\s+\d+:.*?)(?=Slide\s+\d+:|$)', speaker_notes_text, flags=re.IGNORECASE)
                
                for i in range(1, len(notes_by_slide), 2):
                    if i < len(notes_by_slide):
                        slide_header = notes_by_slide[i]
                        slide_content = notes_by_slide[i+1] if i+1 < len(notes_by_slide) else ""
                        
                        # Extract slide number
                        slide_num_match = re.search(r'Slide\s+(\d+):', slide_header)
                        if slide_num_match:
                            slide_num = int(slide_num_match.group(1))
                            
                            # Add intelligence insights if available
                            if slide_num in presentation_intel:
                                intel = presentation_intel[slide_num]
                                
                                # Add slide type indicator for speaker reference
                                type_indicator = f"\n[{intel['type'].value.title()} Slide]"
                                
                                # Add specific insights for data slides
                                if intel['type'] == SlideType.DATA_VISUAL and intel['insights']:
                                    insights_added = []
                                    if intel['insights'].get('trends'):
                                        insights_added.append(f"• Trend insight: {intel['insights']['trends'][0]}")
                                    if intel['insights'].get('outliers'):
                                        insights_added.append(f"• Notable finding: {intel['insights']['outliers'][0]}")
                                    
                                    if insights_added:
                                        slide_content = slide_content.rstrip() + "\n" + "\n".join(insights_added)
                                
                                enhanced_notes.append(slide_header + type_indicator + slide_content)
                            else:
                                enhanced_notes.append(slide_header + slide_content)
                
                speaker_notes_text = "".join(enhanced_notes)
            except Exception as e:
                logger.warning(f"Intelligence enhancement failed: {e}")
        
        state["results"]["speaker_notes"] = speaker_notes_text
        logger.info("Step 5 (Speaker Notes) completed.")

        state["processing_message"] = "+Generation complete!"
        state["processing_step"] = "done"
        state["ui_controls"]["download_disabled"] = "no"

    except GenerationStoppedError:
         state["processing_message"] = "!Generation stopped by user."
         state["processing_step"] = ""  # Hide tabs on stop
         logger.info("Generation process stopped by user request.")
    except Exception as e:
        error_message = f"-An error occurred during step '{state['processing_step']}': {type(e).__name__}"
        state["processing_message"] = error_message
        state["processing_step"] = ""  # Hide tabs on error
        logger.error(f"Pipeline error at step {state['processing_step']}: {e}", exc_info=True)
    finally:
        state["is_generating"] = False
        state["stop_requested"] = False
        _update_generate_button_state(state)

# Download Handlers
def handle_download_outline(state):
    _download_content_as_docx(state, "outline")

def handle_download_visuals(state):
    _download_content_as_docx(state, "visuals")

def handle_download_briefing(state):
    _download_content_as_docx(state, "briefing_info")

def handle_download_mapping(state):
    _download_content_as_docx(state, "mapping")

def handle_download_notes(state):
    _download_content_as_docx(state, "speaker_notes")